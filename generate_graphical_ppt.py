import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_graphical_presentation():
    prs = Presentation()

    # Define Colors
    DEEP_BLUE = RGBColor(0x1B, 0x3A, 0x57)
    TEAL = RGBColor(0x00, 0x80, 0x80)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4)

    # Paths to generated images
    BASE_PATH = r"C:\Users\mimra\.gemini\antigravity\brain\3644e874-a516-42f5-b74d-0bcead2d0928"
    IMG_CORP = BASE_PATH + r"\corporate_leadership_concept_1773217565579.png"
    IMG_BOT = BASE_PATH + r"\swiftbot_whatsapp_mockup_1773217581906.png"
    IMG_LOGISTICS = BASE_PATH + r"\pharma_logistics_infographic_base_1773217597040.png"

    def add_background_shape(slide, color=DEEP_BLUE):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.width = 0

    def add_header(slide, text, color=WHITE):
        # Manually add a header box
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DEEP_BLUE
        shape.line.width = 0
        
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        tf = tx.text_frame
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_shape(slide, DEEP_BLUE)
    
    try:
        slide.shapes.add_picture(IMG_CORP, Inches(4), 0, height=Inches(7.5))
    except: pass

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4), Inches(2))
    tf = tx.text_frame
    p = tf.add_paragraph()
    p.text = "SWIFTSALES"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p = tf.add_paragraph()
    p.text = "AI-POWERED PHARMA"
    p.font.size = Pt(28)
    p.font.color.rgb = TEAL

    # 2. Mission
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "OUR CORE MISSION")
    
    missions = ["INNOVATION", "ACCESSIBILITY", "RELIABILITY"]
    for i, m in enumerate(missions):
        x = Inches(0.5 + i*3.2)
        y = Inches(2.5)
        diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, Inches(2.8), Inches(2.8))
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = TEAL
        diamond.line.color.rgb = DEEP_BLUE
        
        tx = slide.shapes.add_textbox(x, y + Inches(1.1), Inches(2.8), Inches(0.6))
        p = tx.text_frame.add_paragraph()
        p.text = m
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # 3. Leadership
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    try:
        slide.shapes.add_picture(IMG_CORP, 0, 0, width=Inches(10))
    except: pass
    
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5), Inches(8), Inches(2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DEEP_BLUE
    rect.fill.transparency = 0.3
    
    tx = slide.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(7.5), Inches(1.5))
    tf = tx.text_frame
    p = tf.add_paragraph()
    p.text = "\"Building the future of Medicine through AI.\""
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = WHITE
    
    p = tf.add_paragraph()
    p.text = "Malik Muhammad Ejaz, CEO"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = TEAL

    # 8. SwiftBot Demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "SWIFTBOT: REAL-TIME AI")
    
    try:
        slide.shapes.add_picture(IMG_BOT, Inches(3), Inches(1.5), height=Inches(5.5))
    except: pass
    
    features = [("Instant Q&A", 0.5, 2), ("Fuzzy Search", 0.5, 4), ("Fast Orders", 7.5, 2), ("Local AI", 7.5, 4)]
    for text, x, y in features:
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2), Inches(0.8))
        btn.fill.solid()
        btn.fill.fore_color.rgb = TEAL
        
        tx = slide.shapes.add_textbox(Inches(x), Inches(y+0.2), Inches(2), Inches(0.4))
        p = tx.text_frame.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # 11. Logistics
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    try:
        slide.shapes.add_picture(IMG_LOGISTICS, 0, 0, width=Inches(10))
    except: pass
    
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6), Inches(10), Inches(1.5))
    footer.fill.solid()
    footer.fill.fore_color.rgb = DEEP_BLUE
    
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(1))
    p = tx.text_frame.add_paragraph()
    p.text = "STRATEGIC GROWTH: NATIONAL REACH 2026"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Final Save
    save_path = "SwiftSales_Executive_Graphical.pptx"
    prs.save(save_path)
    print(f"Graphical Presentation saved successfully to {save_path}")

if __name__ == "__main__":
    create_graphical_presentation()

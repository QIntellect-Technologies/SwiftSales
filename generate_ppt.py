import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define Colors (Corporate Palette: Deep Blue and Teal)
    DEEP_BLUE = RGBColor(0x1B, 0x3A, 0x57)
    TEAL = RGBColor(0x00, 0x80, 0x80)
    ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
    LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_to_slide(slide, text, font_size=44):
        title_shape = slide.shapes.title
        title_shape.text = text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = DEEP_BLUE

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "SwiftSales Healthcare"
    subtitle.text = "Empowering Healthcare Through Innovation & Distribution\nPakistan's #1 AI-Powered Distribution Network\n\nMalik Muhammad Ejaz | CEO\nMarch 2026"

    # Style Title Slide
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = DEEP_BLUE

    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = TEAL

    # 2. Mission
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "Our Mission")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "To revolutionize pharmaceutical distribution through AI and digital accessibility."
    p = tf.add_paragraph()
    p.text = "• Seamless access to 2,136+ pharmaceutical products."
    p = tf.add_paragraph()
    p.text = "• Serving patients in remote districts via WhatsApp-first technology."
    p = tf.add_paragraph()
    p.text = "• Creating 24/7 digital clinical support for our partners."

    # 3. CEO Message
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "Leadership & Vision")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Malik Muhammad Ejaz - Founder & CEO"
    p = tf.add_paragraph()
    p.text = '"Our vision is to bridge the gap between medicine and patients using state-of-the-art AI technology. We don\'t just deliver products; we deliver health and reliability."'
    p.font.italic = True
    p = tf.add_paragraph()
    p.text = "• 12+ years of industry excellence."
    p = tf.add_paragraph()
    p.text = "• Commitment to transparent, data-driven partnerships."

    # 4. Core Values
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "Core Values")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Integrity: Honesty in every transaction."
    tf.add_paragraph().text = "Innovation: AI-driven solutions like SwiftBot."
    tf.add_paragraph().text = "Reliability: Guaranteed stock and fast fulfillment."
    tf.add_paragraph().text = "Partnership: Success built on mutual growth."
    tf.add_paragraph().text = "Excellence: Highest standards in pharma distribution."

    # 5. Sole Distributor Advantage
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "The Sole Distributor Advantage")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Direct relationships with 34+ top manufacturers."
    tf.add_paragraph().text = "• Market dominance in key territories."
    tf.add_paragraph().text = "• Competitive pricing and exclusive access."
    tf.add_paragraph().text = "• End-to-end supply chain control."

    # 6. Capabilities
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "Capabilities at a Glance")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Always Available: 2,136+ products in stock."
    tf.add_paragraph().text = "• Rapid Fulfillment: Order to delivery < 24 hours."
    tf.add_paragraph().text = "• Efficient Logistics: Specialized cold-chain support."
    tf.add_paragraph().text = "• Partner Support: Dedicated digital agents."

    # 7. AI Attendance System
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "Innovation: AI-Based Attendance")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Automated AI tracking for workforce efficiency."
    tf.add_paragraph().text = "• Real-time workforce management dashboard."
    tf.add_paragraph().text = "• Precision data-driven decision making."
    tf.add_paragraph().text = "• significant reduction in operational overhead."

    # 8. SwiftBot Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "Introducing SwiftBot")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "The world's first WhatsApp-first Pharma AI."
    tf.add_paragraph().text = "• Powered by local OLLAMA AI (Private & Fast)."
    tf.add_paragraph().text = "• Works on slow 3G/4G connections."
    tf.add_paragraph().text = "• Zero app installation required for users."

    # 9. SwiftBot Demo
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "SwiftBot in Action")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Automated Ordering & Q&A."
    tf.add_paragraph().text = "• Voice-like ordering: 'Add 5 Panadol'."
    tf.add_paragraph().text = "• Instant stock validation in < 20ms."
    tf.add_paragraph().text = "• 24/7 customer support without human delays."

    # 10. Capacity Building
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "Capacity Building & Trainings")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Continuous Team Development Programs."
    tf.add_paragraph().text = "• Partner training on AI Ordering Systems."
    tf.add_paragraph().text = "• Skill-building workshops for digital adoption."
    tf.add_paragraph().text = "• Culture of continuous operational improvement."

    # 11. Future Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "Future Strategy")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Expansion & Technology 2026 Roadmap."
    tf.add_paragraph().text = "• National Market Expansion."
    tf.add_paragraph().text = "• Urdu language support for SwiftBot."
    tf.add_paragraph().text = "• AI-driven drug interaction checkers."

    # 12. Call to Action
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_to_slide(slide, "Partner With Us")
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Let's build the future of pharma distribution together."
    tf.add_paragraph().text = "Contact Malik Muhammad Ejaz"
    tf.add_paragraph().text = "Phone: +92 321 7780623"
    tf.add_paragraph().text = "Email: swiftsales.healthcare@gmail.com"
    tf.add_paragraph().text = "Office: Sardar Colony, Rahim Yar Khan"

    save_path = "SwiftSales_Corporate_Presentation.pptx"
    prs.save(save_path)
    print(f"Presentation saved successfully to {save_path}")

if __name__ == "__main__":
    create_presentation()
